'''
    04/08/2026  Export coil data to pptx
    04/09/2026  Export coil data to SVG
    04/11/2026  Debug Plot
    04/12/2026  Export coil data to PDF
'''
from abc import ABC, abstractmethod
    
import collections 
import collections.abc
import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import numpy as np

if __name__ == "__main__":
    import coil_color
    import coil_util as cu
else:
    from . import coil_color
    from . import coil_util as cu

def get_color(c):
    cc = coil_color.color_collection
    
    if c in cc:
        c_ = cc[c]
    else:
        c_ = cc['k']
        
    return (c_[0], c_[1], c_[2])
    
def Polyline(slide, x, y, lcol, lthk, fcol, closed):
    lcol = get_color(lcol)
    fcol = get_color(fcol)

    free_form = slide.shapes.build_freeform(Inches(x[0]), Inches(y[0]))
    free_form.add_line_segments(
        [(Inches(x1), Inches(y1)) for x1, y1 in zip(x[1:],y[1:])],
        close=closed
    )
    
    line_shape = free_form.convert_to_shape()
    
    if closed == False: 
        line_shape.fill.background()
    else:
        if isinstance(fcol, tuple):
            line_shape.fill.solid()
            line_shape.fill.fore_color.rgb = RGBColor(fcol[0], fcol[1], fcol[2])
            # lcol is None, set lcol as fcol if not the default colot works
            if not isinstance(lcol, tuple):
                line_shape.line.color.rgb = RGBColor(fcol[0], fcol[1], fcol[2])
                line_shape.line.width = Inches(0.001) # dummy thinkness
        else:
            line_shape.fill.background()    
    
    if lcol:
        if fcol and lcol == fcol:
            line_shape.line.color.rgb = RGBColor(fcol[0], fcol[1], fcol[2])
            line_shape.line.width = Inches(lthk)
        else:
            line_shape.line.color.rgb = RGBColor(lcol[0], lcol[1], lcol[2])
            line_shape.line.width = Inches(lthk)
        
    line_shape.shadow.inherit     = False
    line_shape.shadow.blur_radius = 0
    line_shape.shadow.distance    = 0
 
    
def draw_transition_ellipse(dev, co):
    ex, ey, ts = co.c1x, co.c1y, co.t_star
    aa, bb = co.axlen, co.bxlen
    xt = aa*np.cos(ts)
    yt = bb*np.sin(ts)
    dd = np.linspace(0, 2*np.pi, 100)
    
    px, py = co.P[2], co.P[3]
    xx1 = co.c1x+aa*np.cos(dd)
    yy1 = co.c1y+bb*np.sin(dd)
    
    xx2 = xx1+co.p_dist
    yy2 = yy1
    
    xs_c = px + co.r_fillet*np.cos(dd)
    ys_c = py + co.r_fillet*np.sin(dd)
    dev.polyline(xs_c, ys_c, lcol="gold", lthk=0.01)
    
    if co.coil_type == cu._coil_type_ellipse_curvature or \
       co.coil_type == cu._coil_type_ellipse_shape:
        xx3 = co.xc_s+co.a_s*np.cos(dd)
        yy3 = co.yc_s+co.b_s*np.sin(dd)
        dev.polyline(xx3, yy3, lcol="blueviolet", lthk=0.01)

    p1x, p1y = ex+xt, ey+yt
    p2x, p2y = px-(ex+xt-px), ey+yt
    dev.polyline([px, p1x], [py, p1y], lcol="cadetblue", lthk=0.01)
    dev.polyline([px, p2x], [py, p2y], lcol="cadetblue", lthk=0.01)
        
class Device():
    def __init__(self, xx, yy):
        self.xmin, self.xmax = min(xx), max(xx)
        self.ymin, self.ymax = min(yy), max(yy)
        self.max_range = max(self.xmax-self.xmin, self.ymax-self.ymin)
        
    @abstractmethod
    def xs_(self, xx): pass
   
    @abstractmethod
    def ys_(self, yy): pass

    @abstractmethod
    def close(self): pass
    
    @abstractmethod
    def polyline(self, xx, yy, lcol, lthk): pass
    
class DevicePPT(Device):
    def __init__(self, fname, xx, yy):
        super().__init__(xx,yy)
        self.fname = fname
        self.ppt = pptx.Presentation()
        self.blank_slide_layout = self.ppt.slide_layouts[6]
        self.slide = self.ppt.slides.add_slide(self.blank_slide_layout)
        
        # Get width and height in EMUs
        width_emu = self.ppt.slide_width
        height_emu = self.ppt.slide_height
        
        # Convert to inches for readability
        width_inches = width_emu / 914400
        height_inches = height_emu / 914400
        min_edge = min(width_inches, height_inches)
        self.scale = min_edge/self.max_range
        
    def xs_(self, xx): return (xx-self.xmin)*self.scale
    def ys_(self, yy): return (yy-self.ymin)*self.scale
    
    def polyline(self, xx, yy, lcol, lthk):
        Polyline(self.slide, self.xs_(xx), self.ys_(yy), lcol, lthk, 'w', False)
    
    def close(self):
        self.ppt.save(self.fname)
     

def save_ppt(coil, fname, trace=False, lcol='b', lthk=0.01, debug=False, lead_l=0, lead_r=0):
    fc = 'w'
    if trace:
        c1, c2, c3 = 'r', 'g', 'b'
        xx, yy, seg = coil.create_geom(trace, lead_l, lead_r)
        dev = DevicePPT(fname, xx, yy)
        
        bit=True
        for i,s in enumerate(seg):
            ii=i+1
            if ii%2 == 0:
                Polyline(dev.slide, dev.xs_(s[0]), dev.ys_(s[1]), c2, lthk, fc, False)
            else:
                Polyline(dev.slide, dev.xs_(s[0]), dev.ys_(s[1]), c1 if bit else c3, lthk, fc, False)
                bit = not bit
    else:
        xx, yy = coil.create_geom(trace, lead_l, lead_r)
        dev = DevicePPT(fname, xx, yy)
        Polyline(dev.slide, dev.xs_(xx), dev.ys_(yy), lcol, lthk, fc, False)
        
    if debug:
        draw_transition_ellipse(dev, coil)
        
    dev.close()
    
'''
    Device SVG
'''
    
_line_format_begin = "<line x1=\"%3.3f\" y1=\"%3.3f\" x2=\"%3.3f\" y2=\"%3.3f\" "
_line_format_end = " style=\"fill:none;stroke:rgb(%d,%d,%d);stroke-width:%3.3f\" />\n"
_polygon_format_end = "style=\"stroke:rgb(%d,%d,%d);stroke-width:%d;fill:rgb(%d,%d,%d);\"/>\n"
_polygon_format_end_nostroke = "style=\"stroke:none;fill:rgb(%d,%d,%d);\"/>\n"
_polygon_format_end_nofill = "style=\"stroke:rgb(%d,%d,%d);stroke-width:%d;fill:none;\"/>\n"
_next_line = "\n"              
_points_per_line = 5
_move_to_cmd = 'M'
_line_to_cmd = "L"
               
class DeviceSVG(Device):
    def __init__(self, fname, xx, yy, hgt=300, ar=1.33, xgap=3, ygap=3):
        super().__init__(xx, yy)
        wid = int(hgt*ar)
        self.fp = open(fname, "wt")
        self.fp.write("<svg version=\"1.1\"\n"\
                      "width=\"%d\" height=\"%d\"\n"
                      "xmlns=\"http://www.w3.org/2000/svg\">\n"\
                      %(int(wid),int(hgt)))
        self.scale = min(hgt,wid)/self.max_range
        self.xx = xx
        self.yy = yy
        self.xgap = xgap
        self.ygap = ygap
        
    def polyline(self, xx, yy, lcol, lthk):
        lc = get_color(lcol)
        lt = self.scale*lthk
        self.fp.write("<polyline points=\"\n")
        self.create_pnt_list(xx, yy)
        self.fp.write(_line_format_end%(\
                      lc[0], 
                      lc[1], 
                      lc[2],
                      1 if lt < 0.001 else lt))
                      
    def create_pnt_list(self, xx, yy):
        xs = self.xs_(xx)
        ys = self.ys_(yy)
        for i, (x1, y1) in enumerate(zip(xs, ys)):
            self.fp.write("%3.3f %3.3f, "%(x1, y1))
            if (i+1)%_points_per_line == 0:
                self.fp.write(_next_line)
        self.fp.write("\"\n")      
        
    def xs_(self, xx): return self.xgap+(xx-self.xmin)*self.scale
    def ys_(self, yy): return self.ygap+(yy-self.ymin)*self.scale
    
    def close(self):
        self.fp.write("</svg>")
        self.fp.close()     

def save_svg(coil, fname, trace=False, lcol='b', lthk=0.01, debug=False, lead_l=0, lead_r=0):
    
    if trace:
        c1, c2, c3 = 'r', 'g', 'b'
        xx, yy, seg = coil.create_geom(trace, lead_l, lead_r)
        dev = DeviceSVG(fname, xx, yy)
        
        bit=True
        for i,s in enumerate(seg):
            ii=i+1
            if ii%2 == 0:
                dev.polyline(s[0], s[1], c2, lthk)
            else:
                dev.polyline(s[0], s[1], c1 if bit else c3, lthk)
                bit = not bit
    else:
        xx, yy = coil.create_geom(trace, lead_l, lead_r)
        dev = DeviceSVG(fname, xx, yy)
        dev.polyline(xx, yy, lcol, lthk)
        
    if debug:
        draw_transition_ellipse(dev, coil)
        
    dev.close()        
    
'''
    PDF Driver
'''    

import zlib
import math

_pdf_header = "%PDF-1.7\n"
_points_inch = 72
_default_nobj = 3

_translate = lambda x,y : "1.0000 0.0000 "\
                          "0.0000 1.0000 "\
                          "%3.4f %3.4f cm\n"%(x,y)
                          
_rotate = lambda phi : "%3.4f %3.4f "\
                       "%3.4f %3.4f "\
                       "0.0000 0.0000 cm\n"%(
                        np.cos(phi), np.sin(phi),
                       -np.sin(phi), np.cos(phi)
                       )
                       
_y_inverse = lambda hgt : "1.0000 0.0000 "\
                          "0.0000 -1.0000 "\
                          "0.0000 %3.4f cm\n"%(
                          hgt
                          )

_x_inverse = lambda wid : "-1.0000 0.0000 "\
                          "0.0000 1.0000 "\
                          "%3.4f 0.0000 cm\n"%(
                          wid
                          )
def color_normalize(c):
    return (c[0]/255., c[1]/255., c[2]/255.)
    
class DevicePDF(Device):
    def __init__(
            self, 
            fname, xx, yy, size=(8.5,11.0), xgap=0.5, ygap=0.5,
            layout_dir='p',
            compression=False):
            
        super().__init__(xx, yy)
        wid = (size[0]-2*xgap)
        hgt = (size[1]-2*ygap)
        self.scale = min(hgt,wid)/self.max_range
        self.xgap = xgap
        self.ygap = ygap
        
        self.obj_list = []
        self.file_size = 0
        self.layout_dir = layout_dir
        self.compression = compression
        self.obj_length = 0
        self.cur_obj_index = 0
        self.rotate = 0
        self.obj_length = 0
        self._sx = 0
        self._sy = 0
        self._ex = size[0]*_points_inch
        self._ey = size[1]*_points_inch

        if layout_dir.upper == 'P':
            self.rotate = 90
            obj_buffer_list = [ _translate(self._ex, 0), 
                                _rotate(util.deg_to_rad(90)),
                                _y_inverse(self._ex)]
        else:
            obj_buffer_list = [ _y_inverse(self._ey)]

        obj_buffer = ''.join(obj_buffer_list) if self.compression else\
                     bytes(''.join(obj_buffer_list), 'utf-8')
        self.obj_list.append(obj_buffer)
        self.obj_length += len(obj_buffer)
        self.fp = open(fname, "wb")
        
    def polyline(self, xx, yy, lcol, lthk, fcol=None, closed=False):
        xs = self.xs_(xx)
        ys = self.ys_(yy)
        lc = color_normalize(get_color(lcol)) if lcol else lcol
        fc = color_normalize(get_color(fcol)) if fcol else fcol
        lt = lthk*self.scale*_points_inch
        
        obj_buffer_list = ["q\n"] #saveDC
        
        if lcol:
            obj_buffer_list.append("%1.4f %1.4f %1.4f RG\n"%(lc[0], lc[1], lc[2]))
            obj_buffer_list.append("%3.3f w\n"%lt)
        
        if fcol:
            obj_buffer_list.append("%1.4f %1.4f %1.4f rg\n"%(fc[0], fc[1], fc[2]))
            
        obj_buffer_list.append("%3.3f %3.3f m\n"%(xs[0],ys[0]))
        
        for x1, y1 in zip(xs[1:],ys[1:]):
            obj_buffer_list.append("%3.3f %3.3f l\n"%(x1,y1))
        
        if closed:
            if lcol and fcol:
                obj_buffer_list.append("b\nQ\n") # close, fill, stroke and restore DC
            elif not isinstance(lcol, color.Color) and fcol:
                obj_buffer_list.append("f\nQ\n") # close, fill, and restore DC
            else:
                obj_buffer_list.append("s\nQ\n") # close, stroke and restore DC
        else:
            obj_buffer_list.append("S\nQ\n")     # stroke and restoreDC
            
        obj_buffer = ''.join(obj_buffer_list) if self.compression else\
                     bytes(''.join(obj_buffer_list), 'utf-8')
        self.obj_list.append(obj_buffer)
        self.obj_length += len(obj_buffer)
        
    def polygon(self, xx, yy, lcol, lthk, fcol):
        self.polyline(xx,yy,lcol,lthk,fcol,True)
    
    #---------------------------------------------------------------
    # Currently, the total number of ojb is 4. 
    # The 4th obj is the ploting commands. 
    #---------------------------------------------------------------
    
    def close(self):

        obj1 = "1 0 obj\n<< /Type /Catalog /Pages 2 0 R>>\nendobj\n"
        obj2 = "2 0 obj\n<< /Type /Pages /Kids [3 0 R] /Count 1>>\nendobj\n"
        obj3 = "3 0 obj\n<<\n/Type /Page\n/Parent 2 0 R\n"\
                "/MediaBox [%3.4f %3.4f %3.4f %3.4f]\n"\
                "/Rotate %d\n"\
                "/Contents 4 0 R\n>>\nendobj\n"%\
                (self._sx, self._sy, self._ex, self._ey, self.rotate)
        
        # Write PDF Header
        self.file_size = 0
        self.fp.write(bytes(_pdf_header,'utf-8'))
        self.file_size += len(_pdf_header)
        obj_pos = [self.file_size]   
        
        # Write Obj 1
        self.fp.write(bytes(obj1,'utf-8'))
        self.file_size += len(obj1)
        obj_pos.append(self.file_size)
        
        # Write Obj 2
        self.fp.write(bytes(obj2,'utf-8'))
        self.file_size += len(obj2)
        obj_pos.append(self.file_size)
        
        # Write Obj 3
        self.fp.write(bytes(obj3,'utf-8'))
        self.file_size += len(obj3)
        obj_pos.append(self.file_size)        

        # Write Obj 4
        if self.compression:
            zip_obj = zlib.compress(bytes(''.join(self.obj_list), 'utf-8'))
            self.obj_length = len(zip_obj)
        
        obj4 = bytes("4 0 obj\n<</Length %d %s>>\nstream\n"%(
                     self.obj_length,
                     "/Filter [/FlateDecode]" if self.compression else ""), 'utf-8')
        self.file_size += len(obj4)
        self.fp.write(obj4)

        if self.compression:
            self.fp.write(zip_obj)
        else:
            for o in self.obj_list:
                self.fp.write(o)
        self.file_size += self.obj_length
        
        obj4 = bytes("\nendstream\nendobj\n",'utf-8')
        self.file_size += len(obj4)
        self.fp.write(obj4)   
        obj_pos.append(self.file_size)
        
        start_xref = self.file_size
        nobj = len(obj_pos)+1
        self.fp.write(bytes("xref\n0 %d\n0000000000 65535 f\n"%nobj,'utf-8'))
        for v in obj_pos:
            self.fp.write(bytes("%010d 00000 n\n"%(v),'utf-8'))
            
        total_nobj = _default_nobj+1
        self.fp.write(bytes("trailer<</Size %d/Root 1 0 R>>\n"%total_nobj,'utf-8'))
        self.fp.write(bytes("startxref\n%d\n"%start_xref,'utf-8'))
        self.fp.write(bytes("%%EOF",'utf-8'))
        self.fp.close()

    def xs_(self, xx): return (self.xgap+(xx-self.xmin)*self.scale)*_points_inch
    def ys_(self, yy): return (self.ygap+(yy-self.ymin)*self.scale)*_points_inch
          
def save_pdf(coil, fname, trace=False, lcol='b', lthk=0.01, debug=False, lead_l=0, lead_r=0):
    
    if trace:
        c1, c2, c3 = 'r', 'g', 'b'
        xx, yy, seg = coil.create_geom(trace, lead_l, lead_r)
        dev = DevicePDF(fname, xx, yy)
        
        bit=True
        for i,s in enumerate(seg):
            ii=i+1
            if ii%2 == 0:
                dev.polyline(s[0], s[1], c2, lthk)
            else:
                dev.polyline(s[0], s[1], c1 if bit else c3, lthk)
                bit = not bit
    else:
        xx, yy = coil.create_geom(trace, lead_l, lead_r)
        dev = DevicePDF(fname, xx, yy)
        dev.polyline(xx, yy, lcol, lthk)
        
    if debug:
        draw_transition_ellipse(dev, coil)
        
    dev.close()                

